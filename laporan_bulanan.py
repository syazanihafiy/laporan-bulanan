# Author: Syazani Hafiy

import os
import psycopg2
import pandas as pd
import win32com.client

from pptx.util import Pt
from pptx.util import Inches
from pptx import Presentation
from configparser import ConfigParser

def config(filename='database.ini', section='postgresql'):
    # create a parser
    parser = ConfigParser()
    # read config file
    parser.read(filename)

    # get section, default to postgresql
    db = {}
    if parser.has_section(section):
        params = parser.items(section)
        for param in params:
            db[param[0]] = param[1]
    else:
        raise Exception('Section {0} not found in the {1} file'.format(section, filename))
    return db

def connect():
    """ Connect to the PostgreSQL database server """
    conn = None
    try:
        # read connection parameters
        params = config()

        # connect to the PostgreSQL server
        print('Connecting to the PostgreSQL database...')
        conn = psycopg2.connect(**params)
		
        # create a cursor
        cur = conn.cursor()
        
	    # execute a statement
        print('PostgreSQL database version:')
        cur.execute('SELECT version()')

        # display the PostgreSQL database server version
        db_version = cur.fetchone()
        print(db_version)
       
	    # close the communication with the PostgreSQL
        cur.close()
    except (Exception, psycopg2.DatabaseError) as error:
        print(error)
        conn.close()
        print('Database connection closed.')

    return conn

def get_laporan(conn):
    # create a cursor
    cur = conn.cursor()
    try:
        
	    # execute a statement
        cur.execute('SELECT * FROM laporan_perbelanjaan')

        results = cur.fetchall()

	    # close the communication with the PostgreSQL
        cur.close()

        column_names = ['id', 'inisiatif', 'agensi', 'perbelanjaan']

        # convert list to dataframe
        df = pd.DataFrame(results, columns=column_names)

    except (Exception, ValueError) as error:
        print(error)

    return df

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

if __name__ == '__main__':
    conn = connect()
    results = get_laporan(conn)

    # generate a powerpoint
    laporan_ppt = Presentation()

    # customize powerpoint format
    laporan_ppt.slide_width = Inches(16)
    laporan_ppt.slide_height = Inches(9)
    slide_layout = laporan_ppt.slide_layouts[6]

    i = 0

    # as long as there are entries, new slide will be generated and each slide will be showing 5 entries each in a single table.
    while (not results.empty):
        # get first 5 entries
        laporan_result = results.reset_index(drop = True)[:5]
        laporan_result = laporan_result[['id', 'inisiatif', 'agensi', 'perbelanjaan']]

        # drop first 5 entries
        results = results.iloc[5:]

        # customize slide format
        slide = laporan_ppt.slides.add_slide(slide_layout)
        slide = laporan_ppt.slides[i]

        # create table
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(13.5), Inches(4)
        shape = slide.shapes.add_table(6, 4, x, y, cx, cy)
        table = shape.table
        table.cell(0,0).text = "#"
        table.cell(0,1).text = "Inisiatif"
        table.cell(0,2).text = "Agensi"
        table.cell(0,3).text = "Perbelanjaan (RM juta)"

        i = i + 1

        # get table from the current slide
        shapes = slide.shapes
        table_list = []
        for shape_idx in range(len(shapes)):
            shape = shapes[shape_idx]
            if shape.shape_type == 19:
                table_list.append(shape_idx)
        
        # insert entries into table
        table = shapes[shape_idx].table
        for j in range(5):
            for k in range(4):
                cell = table.cell(j+1, k)
                cell.text = str(laporan_result.iloc[j, k])

        # customize font size and font family
        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
                    run.font.name = "Calibri"

    # added text title with ( current slide / total slide )
    for l in range (0, i):
        slide = laporan_ppt.slides[l]
        text_top = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(0), Inches(0))
        text_top.text_frame.paragraphs[0]
        p = text_top.text_frame.add_paragraph()
        p.text = "Kementerian Kewangan (" + str(l+1) + "/" + str(i) + ")"
        p.font.name = "Calibri"
        p.font.size = Pt(18)

    # close database connection
    if conn is not None:
        conn.close()
        print('Database connection closed.')

    # generate powerpoint
    laporan_ppt.save('Laporan_Bulanan.pptx')